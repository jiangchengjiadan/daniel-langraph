# src/parser/pptx_parser.py
"""PPT text extraction using python-pptx"""

from pathlib import Path
from typing import List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_parser import BaseParser
from .image_handler import ImageHandler
from ..models import SlideContent, ImageInfo
from ..logging import log
from ..config import config


class PptxParser(BaseParser):
    """Parser for extracting text from PowerPoint files"""

    def __init__(self):
        self.log = log.bind(module="pptx_parser")
        self.image_handler = ImageHandler()

    def get_supported_extensions(self) -> List[str]:
        return ['.pptx', '.ppt']

    def parse(self, file_path: str) -> Tuple[List[SlideContent], List[ImageInfo]]:
        """
        Parse PPTX file to extract text content and images

        Args:
            file_path: Path to the PPTX file

        Returns:
            Tuple of (slides_content, images)
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"PPTX file not found: {file_path}")

        self.log.info(f"Parsing PPTX: {file_path}")

        # Extract text content
        slides_content = self.extract(str(file_path))

        # Extract images
        images_dir = config.images_dir
        images = self.image_handler.extract(
            str(file_path),
            str(images_dir),
            file_path.stem,
        )

        self.log.info(f"Extracted {len(slides_content)} slides and {len(images)} images")
        return slides_content, images

    def extract(self, pptx_path: str) -> List[SlideContent]:
        """
        Extract all text content from a PPTX file

        Args:
            pptx_path: Path to the PPTX file

        Returns:
            List of SlideContent objects
        """
        self.log.info(f"Extracting text from: {pptx_path}")
        prs = Presentation(pptx_path)
        slides_content = []

        for page_number, slide in enumerate(prs.slides, start=1):
            slide_content = self._extract_slide(slide, page_number)
            slides_content.append(slide_content)

        self.log.info(f"Extracted {len(slides_content)} slides")
        return slides_content

    def _extract_slide(self, slide, page_number: int) -> SlideContent:
        """Extract content from a single slide"""
        title = self._extract_title(slide)
        text = self._extract_body_text(slide)
        notes = self._extract_notes(slide)

        return SlideContent(
            page_number=page_number,
            title=title,
            text=text,
            notes=notes,
        )

    def _extract_title(self, slide) -> Optional[str]:
        """Extract title from slide"""
        for shape in slide.shapes:
            if hasattr(shape, "has_title") and shape.has_title:
                title = shape.title.text.strip()
                if title:
                    return title
        return None

    def _extract_body_text(self, slide) -> str:
        """Extract all body text from slide"""
        texts = []
        for shape in slide.shapes:
            # Skip title as it's handled separately
            if hasattr(shape, "has_title") and shape.has_title:
                continue

            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and text not in texts:
                    texts.append(text)

        return "\n\n".join(texts)

    def _extract_notes(self, slide) -> str:
        """Extract speaker notes from slide"""
        notes_slide = slide.notes_slide
        if notes_slide:
            notes_text_frame = notes_slide.notes_text_frame
            if notes_text_frame and notes_text_frame.text:
                return notes_text_frame.text.strip()
        return ""


def extract_text_from_pptx(pptx_path: str) -> List[SlideContent]:
    """
    Convenience function to extract text from PPTX

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        List of SlideContent objects
    """
    parser = PptxParser()
    return parser.extract(pptx_path)
