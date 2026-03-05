# src/parser/image_handler.py
"""Image extraction from PowerPoint files"""

import os
import subprocess
import tempfile
import hashlib
from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models import ImageInfo
from ..logging import log


class ImageHandler:
    """Handler for extracting images from PowerPoint files"""

    def __init__(self):
        self.log = log.bind(module="image_handler")

    def _get_file_hash(self, file_name: str) -> str:
        """Get 8-character hash prefix for a file name"""
        if not file_name:
            return "unknown"
        hash_val = hashlib.md5(file_name.encode("utf-8")).hexdigest()[:8]
        return hash_val

    def extract(self, pptx_path: str, output_dir: str, file_name: str = "") -> List[ImageInfo]:
        """
        Extract all images from a PPTX file

        Args:
            pptx_path: Path to the PPTX file
            output_dir: Directory to save images
            file_name: Optional file name prefix for organizing images

        Returns:
            List of ImageInfo objects
        """
        pptx_path = Path(pptx_path)
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        # Generate hash prefix for this file
        file_hash = self._get_file_hash(file_name or pptx_path.stem)

        self.log.info(f"Extracting images from: {pptx_path}")
        prs = Presentation(pptx_path)
        images_info = []

        for page_number, slide in enumerate(prs.slides, start=1):
            page_images = self._extract_slide_images(slide, page_number, output_dir, file_hash)
            images_info.extend(page_images)

        self.log.info(f"Extracted {len(images_info)} images")
        return images_info

    def _extract_slide_images(
        self, slide, page_number: int, output_dir: Path, file_hash: str
    ) -> List[ImageInfo]:
        """Extract images from a single slide"""
        images = []
        image_idx = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_idx += 1
                image_info = self._save_image(
                    shape, page_number, image_idx, output_dir, file_hash
                )
                if image_info:
                    images.append(image_info)

        return images

    def _is_emf_or_wmf(self, ext: str, blob: bytes) -> bool:
        """Check if the image is EMF or WMF format"""
        if ext.lower() in ["emf", "wmf"]:
            return True
        # Check magic bytes for EMF/WMF
        if len(blob) >= 4:
            # WMF magic: 0xD7CDC69A (little endian) or 0x9AC6CDD7 (big endian)
            magic = blob[:4]
            # Check for WMF magic number (Aldus placeable metafile)
            if magic == b'\xd7\xcd\xc6\x9a' or magic == b'\xd7\xcd\xc6\x9a'[::-1]:
                return True
            # Check for EMF magic (Windows GDI enhanced metafile)
            if blob[0:4] == b'EMF\x00':
                return True
        return False

    def _convert_blob_to_png(self, blob: bytes, output_path: Path) -> bool:
        """
        Convert EMF/WMF blob to PNG using LibreOffice + pdftoppm

        Args:
            blob: Raw image bytes
            output_path: Path for output PNG file

        Returns:
            True if conversion successful
        """
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                tmpdir = Path(tmpdir)

                # Save blob to temp file
                ext = "emf"  # EMF is the base format
                temp_emf = tmpdir / f"image.{ext}"
                temp_emf.write_bytes(blob)

                # Step 1: Convert EMF/WMF to PDF using LibreOffice
                pdf_path = tmpdir / "image.pdf"
                cmd = [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(tmpdir),
                    str(temp_emf),
                ]
                result = subprocess.run(cmd, capture_output=True, text=True)

                if result.returncode != 0:
                    self.log.warning(f"LibreOffice conversion failed: {result.stderr}")
                    return False

                if not pdf_path.exists():
                    self.log.warning(f"PDF not created for EMF")
                    return False

                # Step 2: Convert PDF to PNG using pdftoppm
                ppm_path = tmpdir / "image"
                cmd = [
                    "pdftoppm",
                    "-png",
                    "-f",
                    "1",
                    "-l",
                    "1",
                    str(pdf_path),
                    str(ppm_path),
                ]
                subprocess.run(cmd, capture_output=True)

                # Find the generated PNG file (pdftoppm adds -1 suffix)
                png_files = list(tmpdir.glob("*-1.png"))
                if png_files:
                    png_files[0].rename(output_path)
                    return True

                return False

        except Exception as e:
            self.log.warning(f"Failed to convert EMF to PNG: {e}")
            return False

    def _save_image(
        self, shape, page_number: int, image_idx: int, output_dir: Path, file_hash: str
    ) -> Optional[ImageInfo]:
        """Save image and return ImageInfo"""
        try:
            image = shape.image
            blob = image.blob
            ext = image.ext.lower() if image.ext else "bin"

            # Build filename with hash prefix: {hash}_{page}_{idx}.{ext}
            base_name = f"{file_hash}_{page_number}_{image_idx}"

            # Check if it's EMF or WMF
            is_vector = self._is_emf_or_wmf(ext, blob)

            if is_vector:
                # Convert EMF/WMF to PNG
                output_path = output_dir / f"{base_name}.png"
                if self._convert_blob_to_png(blob, output_path):
                    return ImageInfo(
                        page_number=page_number,
                        image_idx=image_idx,
                        path=str(output_path),
                        mimetype="image/png",
                    )
                else:
                    # Fallback: save as-is with .emf extension
                    output_path = output_dir / f"{base_name}.emf"
                    output_path.write_bytes(blob)
                    return ImageInfo(
                        page_number=page_number,
                        image_idx=image_idx,
                        path=str(output_path),
                        mimetype="image/emf",
                    )
            else:
                # Standard image format
                if ext not in ["png", "jpg", "jpeg", "gif", "webp", "bmp"]:
                    ext = "bin"

                output_path = output_dir / f"{base_name}.{ext}"
                output_path.write_bytes(blob)

                return ImageInfo(
                    page_number=page_number,
                    image_idx=image_idx,
                    path=str(output_path),
                    mimetype=f"image/{ext}",
                )

        except Exception as e:
            self.log.warning(f"Failed to extract image from shape: {e}")
            return None


def extract_images(pptx_path: str, output_dir: str, file_name: str = "") -> List[ImageInfo]:
    """
    Convenience function to extract images from PPTX

    Args:
        pptx_path: Path to the PPTX file
        output_dir: Directory to save images
        file_name: Optional file name prefix

    Returns:
        List of ImageInfo objects
    """
    handler = ImageHandler()
    return handler.extract(pptx_path, output_dir, file_name)
